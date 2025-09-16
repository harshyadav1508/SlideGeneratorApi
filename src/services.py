# src/services.py
import google.generativeai as genai
import json
import re
import logging  # --- LOGGING: Import the logging module ---
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .config import GEMINI_API_KEY

# --- LOGGING: Get a logger instance for this module ---
logger = logging.getLogger(__name__)

# Configure the Gemini API client
genai.configure(api_key=GEMINI_API_KEY)


class PresentationStyles:
    TITLE_FONT_NAME = 'Calibri Light'
    BODY_FONT_NAME = 'Calibri'
    TITLE_FONT_SIZE = Pt(36)
    BODY_FONT_SIZE = Pt(18)
    TITLE_COLOR = RGBColor(0x0A, 0x2C, 0x52)  # Dark Blue
    BODY_COLOR = RGBColor(0x33, 0x33, 0x33)  # Dark Gray


class LLMService:
    """
    Service for interacting with the Google Gemini LLM.
    """

    def __init__(self, model_name="gemini-2.0-flash"):
        self.model = genai.GenerativeModel(model_name)
        self.cache = {}
        self.max_cache_size = 100

    async def generate_slides_content(self, topic: str, num_slides: int) -> dict:
        """
        Generates structured presentation content using the LLM, with caching.
        """
        cache_key = (topic.lower(), num_slides)

        if cache_key in self.cache:
            logger.info(f"CACHE HIT for topic: '{topic}'")
            return self.cache[cache_key]

        logger.info(f"CACHE MISS for topic: '{topic}'. Calling Gemini API.")

        # --- FIX: Added a stronger instruction to ensure valid JSON syntax ---
        prompt = f"""
        Generate a presentation about "{topic}".
        The presentation should have approximately {num_slides} slides.

        Please provide the output in a single, valid JSON object. Do not include any text, introductory phrases, or markdown formatting like ```json before or after the JSON object.
        The entire response should be only the JSON content.
        VERY IMPORTANT: Ensure the JSON is perfectly formatted. Pay close attention to syntax, especially commas between objects in lists.

        IMPORTANT CONTENT RULES:
        - Each bullet point or definition should be descriptive and detailed, ideally between 15 and 20 words long.
        - You can use markdown for emphasis. Use **double asterisks** for bold and __double underscores__ for underline.
        - For bullet points, provide them as a list of strings without any leading characters like '*'.

        The JSON object must have a single key "slides", which is a list of slide objects.
        Each slide object must have two keys: "layout" and "content".

        The "layout" key must be one of the following three strings:
        1. "title_slide"
        2. "bullet_points"
        3. "two_column"

        The "content" key must be an object with the following structure based on layout:
        - For "title_slide": {{"title": "...", "subtitle": "..."}}
        - For "bullet_points": {{"title": "...", "points": ["...", "..."]}}
        - For "two_column": {{"title": "...", "left_column": {{"heading": "...", "points": ["...", "..."]}}, "right_column": {{"heading": "...", "points": ["...", "..."]}} }}

        Also, please include a slide with source citations at the end, using the "bullet_points" layout.

        Now, generate the complete JSON for the presentation on "{topic}".
        """

        try:
            response = await self.model.generate_content_async(prompt)
            response_text = response.text.strip()
            if response_text.startswith("```json"):
                response_text = response_text[7:]
            if response_text.endswith("```"):
                response_text = response_text[:-3]

            parsed_json = json.loads(response_text)

            self.cache[cache_key] = parsed_json

            if len(self.cache) > self.max_cache_size:
                oldest_key = next(iter(self.cache))
                del self.cache[oldest_key]
                logger.info("Cache max size reached. Removed oldest item.")

            return parsed_json
        except (json.JSONDecodeError, Exception) as e:
            # --- LOGGING: Replaced print with logger.error and added more detail ---
            logger.error(f"Error decoding LLM response: {e}", exc_info=True)
            raw_response = "Response object not available"
            if 'response' in locals() and hasattr(response, 'text'):
                raw_response = response.text
            # Log the problematic raw response for easy debugging
            logger.debug(f"Raw LLM Response that failed parsing:\n{raw_response}")
            return {"error": "Failed to generate or parse content from LLM.", "details": str(e)}


class SlideBuilderService:
    # ... (The rest of this class remains exactly the same) ...
    """
    Service for building a .pptx presentation file.
    """

    def _process_markdown_to_runs(self, paragraph, text):
        paragraph.clear()

        pattern = r'(\*\*|__)(.*?)\1'
        parts = re.split(pattern, text)

        run = paragraph.add_run()
        run.text = parts[0]

        for i in range(1, len(parts), 3):
            delimiter = parts[i]
            content = parts[i + 1]

            run = paragraph.add_run()
            run.text = content

            if delimiter == '**':
                run.font.bold = True
            elif delimiter == '__':
                run.font.underline = True

            if (i + 2) < len(parts):
                run = paragraph.add_run()
                run.text = parts[i + 2]

    def create_presentation(self, presentation_data: dict, output_path: str, aspect_ratio: str) -> None:
        """
        Creates a .pptx file from structured data.
        """
        # --- NEW: Choose the template file based on the aspect ratio ---
        template_path = 'template_16_9.pptx'
        if aspect_ratio == "4:3":
            template_path = 'template_4_3.pptx'

        try:
            prs = Presentation(template_path)
            print(f"Using template: {template_path}")
        except FileNotFoundError:
            print(f"WARNING: '{template_path}' not found. Using default blank presentation.")
            prs = Presentation()


        for slide_info in presentation_data.get("slides", []):
            layout_name = slide_info.get("layout")
            content = slide_info.get("content")

            if layout_name == "title_slide":
                self._add_title_slide(prs, content)
            elif layout_name == "bullet_points":
                self._add_bullet_points_slide(prs, content)
            elif layout_name == "two_column":
                self._add_two_column_slide(prs, content)

        prs.save(output_path)

    def _apply_font_style(self, text_frame, is_title=False):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if is_title:
                    if run.font.bold is not True: run.font.name = PresentationStyles.TITLE_FONT_NAME
                    run.font.size = PresentationStyles.TITLE_FONT_SIZE
                    run.font.color.rgb = PresentationStyles.TITLE_COLOR
                else:
                    if run.font.bold is not True: run.font.name = PresentationStyles.BODY_FONT_NAME
                    run.font.size = PresentationStyles.BODY_FONT_SIZE
                    run.font.color.rgb = PresentationStyles.BODY_COLOR

    def _add_title_slide(self, prs, content):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        self._process_markdown_to_runs(title.text_frame.paragraphs[0], content.get("title", "Presentation Title"))
        self._process_markdown_to_runs(subtitle.text_frame.paragraphs[0], content.get("subtitle", ""))

        self._apply_font_style(title.text_frame, is_title=True)
        self._apply_font_style(subtitle.text_frame)

    def _add_bullet_points_slide(self, prs, content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        self._process_markdown_to_runs(title_shape.text_frame.paragraphs[0], content.get("title", "Slide Title"))
        self._apply_font_style(title_shape.text_frame, is_title=True)

        tf = body_shape.text_frame
        tf.clear()

        for point_text in content.get("points", []):
            p = tf.add_paragraph()
            self._process_markdown_to_runs(p, point_text)
            p.level = 0

        self._apply_font_style(tf)

    def _add_two_column_slide(self, prs, content):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        self._process_markdown_to_runs(title_shape.text_frame.paragraphs[0], content.get("title", "Two Column Title"))
        self._apply_font_style(title_shape.text_frame, is_title=True)

        left_data = content.get("left_column", {})
        left_tf = slide.placeholders[1].text_frame
        left_tf.clear()

        if left_data.get("heading"):
            p_heading = left_tf.add_paragraph()
            self._process_markdown_to_runs(p_heading, left_data["heading"])  # This now handles bolding
            p_heading.level = 0

        for point_text in left_data.get("points", []):
            p_point = left_tf.add_paragraph()
            self._process_markdown_to_runs(p_point, point_text)
            p_point.level = 1

        right_data = content.get("right_column", {})
        right_tf = slide.placeholders[2].text_frame
        right_tf.clear()

        if right_data.get("heading"):
            p_heading = right_tf.add_paragraph()
            self._process_markdown_to_runs(p_heading, right_data["heading"])  # This now handles bolding
            p_heading.level = 0

        for point_text in right_data.get("points", []):
            p_point = right_tf.add_paragraph()
            self._process_markdown_to_runs(p_point, point_text)
            p_point.level = 1

        self._apply_font_style(left_tf)
        self._apply_font_style(right_tf)